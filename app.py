from flask import Flask, request, render_template, send_file, flash, redirect, url_for
from langchain import hub
from langchain_openai import ChatOpenAI
from langchain.agents import AgentExecutor, create_react_agent
import wikipedia
from pptx import Presentation
from dotenv import load_dotenv
from langchain_core.tools import Tool
from langchain.prompts.prompt import PromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
import secrets
import json
import os
import tiktoken
from pydantic import BaseModel, ValidationError


def count_tokens(text, model="gpt-3.5-turbo"):
    encoding = tiktoken.encoding_for_model(model)
    return len(encoding.encode(text))


load_dotenv()

app = Flask(__name__)


class LLMOutput(BaseModel):
    input: str
    summary: str
    ppt_file: str


# Generates a 32-character random hex string
app.secret_key = secrets.token_hex(16)

# Load the ReAct prompt template
react_prompt = hub.pull("hwchase17/react")

# Initialize the OpenAI model
model = ChatOpenAI(
    temperature=0, model_name="gpt-3.5-turbo", stop=["\nObservation", "Observation"]
)

# Define the tools for the agent

# Wikipedia search tool
# This tool will search Wikipedia for the given query and return the summary


def wikipedia_search(query):
    wikipedia.set_lang("en")

    try:
        page = wikipedia.page(query)
        summary = wikipedia.summary(query)
        # Truncate content to a maximum of 1000 characters
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000, chunk_overlap=0)
        truncated_content = text_splitter.split_text(
            summary)[0]  # Take the first chunk
        return truncated_content
    except:
        return "Error"


# Create PowerPoint presentation tool
# This tool will create a PowerPoint presentation from the given content
def create_ppt(content):
    with open("content.txt", "w", encoding="utf-8") as file:
        file.write(content)

    # Ensure content is valid JSON
    if isinstance(content, str):
        # Replace single quotes with double quotes
        content = content.replace("'", '"')
        data = json.loads(content)

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]
    title_placeholder.text = data["title"].capitalize()
    body_placeholder.text = f" Summary: {data['summary']}"

    # Use LangChain's RecursiveCharacterTextSplitter to split content
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=380,  # Maximum number of characters per chunk
        chunk_overlap=0,  # Overlap between chunks to maintain context
        separators=[
            "\n",
            " ",
            "",
        ],  # Prioritize splitting by new lines, then spaces, then characters
    )
    # Split the content into chunks
    content_chunks = text_splitter.split_text(data["content"])
    for chunk in content_chunks:
        slide = prs.slides.add_slide(slide_layout)
        body_placeholder = slide.placeholders[1]
        body_placeholder.text = chunk

    # Save the PowerPoint file in the "powerpoints" folder
    ppt_filename = os.path.join("powerpoints", f"{data['title']}.pptx")
    prs.save(ppt_filename)
    return ppt_filename


tools_of_agent = [
    Tool(
        name="Wikipedia_search",
        func=wikipedia_search,
        description="useful for when you need to search from Wikipedia",
    ),
    Tool(
        name="Create_Powerpoint_presentation",
        func=lambda args: create_ppt(args),  # Pass the full content
        description="useful for when you need to create an PowerPoint presentation",
    ),
]

template = """First, perform a Wikipedia search for "{name_of_topic}" and summarize its content. Ensure that the summary strictly pertains to "{name_of_topic}" and does not deviate to related topics. Provide the output in a JSON format like this:
{{
    "title": "{name_of_topic}",
    "summary": "summarized content",
    "content": "full content"
}}
Then, pass the JSON object to create a PowerPoint presentation file titled "{name_of_topic}". Make sure to include the full content in the presentation. At last, return the output in the below JSON format:
{{
    "input": "{name_of_topic}",
    "summary": "summarized content",
    "ppt_file": "{name_of_topic}.pptx"
}}
"""
prompt_template = PromptTemplate(
    template=template, input_variables=["name_of_topic"])


@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")


@app.route("/search", methods=["POST"])
def search():
    if request.method == "POST":
        topic = request.form["topic"].lower()
        agent = create_react_agent(
            llm=model, tools=tools_of_agent, prompt=react_prompt)
        agent_executor = AgentExecutor(
            agent=agent, tools=tools_of_agent, verbose=True, handle_parsing_errors=True)

        # Format the prompt string
        formatted_prompt = prompt_template.format_prompt(name_of_topic=topic)

        # Convert formatted_prompt to a string if necessary
        if hasattr(formatted_prompt, "to_string"):
            formatted_prompt = formatted_prompt.to_string()
        elif hasattr(formatted_prompt, "text"):
            formatted_prompt = formatted_prompt.text

        # Debug token count
        token_count = count_tokens(formatted_prompt)
        print(f"Token count: {token_count}")

        # Truncate the prompt if it exceeds the token limit
        max_prompt_length = 4000
        if len(formatted_prompt) > max_prompt_length:
            formatted_prompt = formatted_prompt[:max_prompt_length]

        try:
            # Execute the agent with the formatted prompt
            result = agent_executor.invoke(input={"input": formatted_prompt})
            # Debug: Print the agent's output
            print("Agent Output:", result['output'])
            with open("output.txt", "w", encoding="utf-8") as file:
                file.write(result["output"])
            try:
                # Validate the JSON using Pydantic
                validated_output = LLMOutput.model_validate(
                    json.loads(result.get("output")))
                # print("Validated LLM Output:", validated_output)
                flash(
                    f"The PowerPoint file '{validated_output.input}'.pptx has been downloaded successfully!", "success")
            except (ValidationError, ValueError, json.JSONDecodeError) as e:
                flash(f"Invalid LLM output: {e}", "error")
                return redirect(url_for("index"))

            # Download PowerPoint generation
            ppt_filename = f"{validated_output.input}.pptx"
            return send_file(f"./powerpoints/{ppt_filename}", as_attachment=True)
        except Exception as e:
            flash(str(e), "error")
            return redirect(url_for("index"))


if __name__ == "__main__":
    app.run(debug=True)
